# -*- coding: utf-8 -*-
import pandas as pd
from scipy import stats
import scipy as sp
import numpy as np
import pathlib
import xarray as xr
import datetime
import matplotlib.pyplot as plt
import pptx
import io
import statsmodels.api as sm

colors = plt.rcParams['axes.prop_cycle'].by_key()['color']

# def fit_langley(langley):
#     lrg_res = stats.linregress(langley.index, langley)
#     lser=pd.Series({'slope': lrg_res.slope, 'intercept': lrg_res.intercept, 'stderr': lrg_res.stderr})
#     return lser

def fit_langley(langley, weighted = True):
    """
    

    Parameters
    ----------
    langley : TYPE
        DESCRIPTION.
    weighted : bool, optional
        If the fit is to be weighted. This is a good idea due to the change in density of data points as the AMF decreases. The default is False.

    Returns
    -------
    None.

    """
    y = langley
    x = langley.index
    x = sm.add_constant(x)
    
    if not weighted:
        mod = sm.OLS(y,x)
    else:
        idx = langley.index
        wt = idx[:-1] - idx[1:]
        wx = np.arange(wt.shape[0])
        f = sp.interpolate.interp1d(wx, wt, bounds_error=False, fill_value='extrapolate')
        w = np.append(wt, f(wx[-1] + 1))
        mod = sm.WLS(y, x, weights = w)

    res = mod.fit()
    
#     resdict[col] = res

    ### test for curved residual
    # print(type(res.resid))
    # print(res.resid)
    
    x = res.resid.index.values
    y = res.resid.values
    pfres = np.polyfit(x,y,2)    

    lser=pd.Series({'slope': res.params[1], 'intercept': res.params[0], 'slope_stderr': res.bse[1], 'intercept_stderr': res.bse[0], 'resid_curvature': pfres[0]})
    out = {}
    out['res_series'] = lser
    out['res'] = res
    return out

def plot_langley(lang_ds, wls = 'all', date = 'Test'):
    if wls == 'all':
        wls = lang_ds.wavelength.values
    else:
        pass
#     wl = 500
    for wl in wls:
        lan = lang_ds.langleys.to_pandas()[wl]
        resid = lang_ds.langley_fit_residual.to_pandas()[wl]
        f,aa = plt.subplots(2, sharex=True, gridspec_kw={'hspace': 0})
        alan, ares = aa
        
        fitres = lang_ds.langley_fitres.to_pandas().loc[wl]
        fit = fitres.intercept + fitres.slope * lang_ds.airmass.to_pandas()
        fit.plot(ax = alan)
        lan.plot(ax = alan, ls = '', marker = '.')
        g = alan.get_lines()[-1]
        g.set_markersize(2)
        
        fitres['corrmatrixdet'] = lang_ds.langley_residual_correlation_prop.values
        
        alan.text(0.01,0.05, fitres.__str__(), transform=alan.transAxes, fontsize = 'x-small')

        resid.plot(ax = ares)
        ares.grid()
        
        alan.set_title(f'{date} - channel {wl}nm')
    return f,aa

class Langley(object):
    def __init__(self, parent, langleys, langley_fit_settings = None):
        self.parent = parent
        langleys.columns.name = 'wavelength'
        self.langleys = langleys
        self.langley_fit_settings = langley_fit_settings
        self._langley_fitres = None
        self._langley_residual_correlation_prop = None
        self._langley_fit_residual = None
    
    def to_xarray_dataset(self):
        ds = xr.Dataset({'langleys': self.langleys,
                         'langley_fit_residual': self.langley_fit_residual,
                         'langley_fitres': self.langley_fitres,
                         'langley_residual_correlation_prop': self.langley_residual_correlation_prop['determinant'],
                         'sp02_serial_no': self.parent.raw_data.serial_no.values[0]})
        return ds
    
    def save2netcdf(self, path2file):
        ds = self.to_xarray_dataset()
        ds.to_netcdf(path2file)
        return ds
    
    @property
    def langley_fit_residual(self):
        if isinstance(self._langley_fit_residual, type(None)):
            self._fit_langles()
        return self._langley_fit_residual
    
    @property
    def langley_fitres(self):
        if isinstance(self._langley_fitres, type(None)):
            self._fit_langles()
        return self._langley_fitres
    
    @property
    def langley_residual_correlation_prop(self):
        if isinstance(self._langley_residual_correlation_prop, type(None)):
            self._fit_langles()
        return self._langley_residual_correlation_prop
    
    def _fit_langles(self):
        langleys = self.langleys
            
        # df_langres = pd.DataFrame(index=['slope', 'intercept', 'stderr'])
        fit_res_dict = {}
        resid_dict = {}
# !!!!!!!!!!!!!!!!! Baustelle
        for wl in langleys:
            # lrg_res = stats.linregress(langleys.index, langleys[wl])
            # df_langres[wl] = pd.Series({'slope': lrg_res.slope, 'intercept': lrg_res.intercept, 'stderr': lrg_res.stderr})
            out =  fit_langley(langleys[wl])
            # df_langres[wl] = out['res_series']
            fit_res_dict[wl] = out['res_series']
            resid_dict[wl] = out['res'].resid
        
        df_langres = pd.DataFrame(fit_res_dict)
        df_langres = df_langres.transpose()
        df_langres.index.name = 'wavelength'
        df_langres.columns.name = 'fit_results'
        self._langley_fitres = df_langres
        
        df_resid = pd.DataFrame(resid_dict)
        df_resid.columns.name = 'wavelength'
        
        resid_corr_matrix = df_resid.corr()
        lrcp = {}
        lrcp['determinant'] = np.linalg.det(resid_corr_matrix)
        
        self._langley_fit_residual = df_resid
        self._langley_residual_correlation_prop = lrcp
        return df_langres

class Calibration(object):
    def __init__(self, raw2lang_output):
        self.raw2lang_output = raw2lang_output
        self.langley_fitres = raw2lang_output['ds_results']
        
    
#     sp02.calibration.get_best_results#, save_as_result = save)

    def plot_od(self, sn, ax = None, show_molecular = False):
        best_10 = self
        fitres_mean = best_10.langley_fitres.langley_fitres.mean(dim = 'datetime').to_pandas()
        wl = best_10.langley_fitres.channle_wavelengths
        if len(wl.dims) == 1:
            wl = wl.to_pandas()
        else:
            wl = wl[:,0].to_pandas()
        fitres_mean.rename(wl, inplace=True)
        fitres_mean['slope'] = fitres_mean.slope.abs()
    
        a = fitres_mean.slope.plot(ls = '', marker = 'o', ax = ax, label = f'sn{sn}')
        a.set_yscale('log')
        a.set_ylabel('OD')
        a.set_xscale('log')

        if show_molecular:
            x = np.linspace(368, 1050, 5)
            y = 1e10 * x**-4
            g, = plt.plot(x, y)
            g.set_zorder(0)
            g.set_linestyle('--')
            g.set_color('0.6')
            g.set_dash_capstyle('round')
            g.set_label('molecular')
        a.legend()  
        return a
        
    def plot_langley_results(self, intensity4ts = 'resid_curvature', add2pptx = False,test = False, tollerance = [2.5, 2.5]):
        """
        slope 	intercept 	slope_stderr 	intercept_stderr 	resid_curvature

        Parameters
        ----------
        result_ds : TYPE
            DESCRIPTION.
        intensity4ts : TYPE, optional
            DESCRIPTION. The default is ''.
        add2pptx : TYPE, optional
            DESCRIPTION. The default is False.
        test : TYPE, optional
            DESCRIPTION. The default is False.
        tollerance : TYPE, optional
            DESCRIPTION. The default is [2.5, 2.5].

        Returns
        -------
        out : TYPE
            DESCRIPTION.

        """
        ds = self.langley_fitres
        out = {}
        if add2pptx:
            date = pd.to_datetime(datetime.datetime.now()).date()
            date = str(date)
            date = date.replace('-','')
            path2ppt = f'/mnt/telg/projects/sp02/calibration/langley_calibration_summary_{date}.ppt'
            path2ppt
            prs = pptx.Presentation()
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]

            title.text = "SP02 revival"
            now = pd.to_datetime(datetime.datetime.now())

            subtitle.text = f"created on {now.date()} by Hagen"



        plt.rcParams['hatch.linewidth'] = 4
        for wl in ds.wavelength.values:
            f,a = plt.subplots()
            res = ds.langley_fitres.sel(wavelength = wl).to_pandas()

            rest = res.copy()

            restdna = rest.dropna()

            cm = plt.cm.plasma_r
            pc = plt.scatter(restdna.slope,restdna.intercept_stderr, c = np.e**restdna.intercept, s = 40, cmap = cm, 
                       )
            cb = f.colorbar(pc)
            cb.set_label('Langley fit - intercept')
            zm = (np.e**restdna.intercept).median()
            zstd  = (np.e**restdna.intercept).std()
            cb.ax.axhline(zm, color = 'black')
            hs = cb.ax.axhspan(zm-zstd, zm+zstd)
            hs.set_facecolor([0,0,0,0])
            hs.set_edgecolor([0,0,0,0.5])
            hatch = hs.set_hatch('//')

            a.set_title(f'channle wavelength: {wl}nm')
            a.set_xlabel('Langley fit - slope (mV)')
            a.set_ylabel('Langley fit - std of residual')

        # langley over time
            f_ivt,a = plt.subplots()
            out['tp_restdna'] = restdna
            restdna.intercept.plot(ax = a)
            pc = a.scatter(restdna.index, np.e**restdna.intercept, c = restdna[intensity4ts], cmap=plt.cm.gnuplot)
            f_ivt.colorbar(pc)
            a.set_xlabel('')
            a.set_ylabel('Langley fit - intercept')
            a.set_title('Langley fit intercept as a function of time')

            if add2pptx:

                blank_slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(blank_slide_layout)
                add_fig2slide(f, slide)
                add_fig2slide(f_ivt, slide, top = 7.5/2)
#             if test:
            break
        if add2pptx:
            prs.save(path2ppt)
        return out
    
    def get_best_results(self, top = 10):
        ds = self.langley_fitres
        def sortbyparam(ds, param, top):
            df = (ds.langley_fitres.sel(linreg_param = f'{param}_stderr') / ds.langley_fitres.sel(linreg_param = param)).mean(dim = 'wavelength').to_pandas()
            df = df[df>0]
            df.sort_values(inplace=True)
            out = {}
            df = df.iloc[:top]
            out = ds.sel(datetime =  df.index.values)
            return out
        out = {}
        slopestdrmin = sortbyparam(ds,'slope', top)
        interceptstdrmin = sortbyparam(ds,'intercept',top)
        curvature_date = ds.langley_fitres.sel(linreg_param = 'resid_curvature').mean(dim = 'wavelength').to_pandas().abs().sort_values().index[:top].values
        curvature = ds.sel(datetime = curvature_date)
        df_cmp = ds.correlation_matrix_properties.to_pandas()
        determimin_dates = df_cmp.sort_values('determinant', ascending=False).index[:top].values
        determimin = ds.sel(datetime = determimin_dates)
    #     uberlap = ds.datetime.values
    #     for other in [interceptstdrmin, 
    #                  curvature, 
    #     #              determimin,
    #                  ]:
    #         uberlap = np.intersect1d(uberlap, other)
    #     out['crossection'] = uberlap

        out['slope_stderr'] = slopestdrmin
        out['intercept_stderr'] = interceptstdrmin
        out['curvature'] = curvature
        out['corr_matrix_determinant_max'] = determimin
        out['ds_results'] = interceptstdrmin
        return Calibration(out)
    
    def plot_good_enough(self, best=None, col = 'A'):
        # col = 'A'
        df = self.langley_fitres.langley_fitres.sel(wavelength = col).to_pandas()
        df_corr = self.langley_fitres.correlation_matrix_properties.to_pandas()
        # ['intercept_stderr'].langley_fitres.sel(wavelength = col).to_pandas()

        df.sort_values('intercept_stderr', inplace=True)
        # df
        df_corr = df_corr.loc[df.index]
        # return df, df_corr  
        
        df.index = range(df.shape[0])
        df_corr.index = range(df.shape[0])
        inter_v0 = np.e**df.intercept
        df['intercept_stderr'] = np.e**(df['intercept'] + df['intercept_stderr']) - inter_v0
        df['intercept'] = inter_v0
        
        

        if not isinstance(best, type(None)):
            df = df.iloc[:best]
            df_corr = df_corr.iloc[:best]

        f, aa = plt.subplots(4, sharex=True, gridspec_kw={'hspace': 0})
        f.set_figheight(f.get_figheight() * 1.5)
        a_int = aa[0]
        a_slope = aa[1]
        a_curv = aa[2]
        a_corr = aa[3]
        
        ### intercept

        df.intercept.plot(ax = a_int, marker = 'o')
        at_int = a_int.twinx()
        df.intercept_stderr.plot(ax = at_int, color = colors[1], marker = 'o')

        ### slope
        df.slope.plot(ax = a_slope, marker = 'o')
        at_slope = a_slope.twinx()
        df.slope_stderr.plot(ax = at_slope, color = colors[1], marker = 'o')

        ### curvature
        df.resid_curvature.plot(ax = a_curv, marker = 'o')

        ### correlation
        df_corr.plot(ax = a_corr, marker = 'o')        

        ### Labels 
        a_int.set_ylabel('V0 (mV)')
        at_int.set_ylabel('stderr', color = colors[1])

        a_slope.set_ylabel('slope')
        at_slope.set_ylabel('stderr', color = colors[1])

        a_curv.set_ylabel('curvature')
        # return a_corr,df_corr
        return aa
    
    def plot_data_with_fit(self, top = 0):
        which = top

        path2ncfld = pathlib.Path('./langleys/')
        best10 = self
#         which = 0

        best10sel = best10.langley_fitres
        dtt = best10sel.langley_fitres.sel(linreg_param = 'intercept_stderr', wavelength = 'A').to_pandas().sort_values().index[which]
        date = pd.to_datetime(dtt)
        date = f'{date.year:04d}{date.month:02d}{date.day:02d}'

        # sn = int(best10sel.sp02_serial_no.values[0])
        sn = best10sel.sp02_serial_no.values
        if len(sn.shape) == 0:
            sn = int(sn)
        else:
            sn = int(sn[0])
        
        fn = f'sn{sn}_{date}_am.nc'

        ds = xr.open_dataset(path2ncfld.joinpath(fn))

        f, aa = plot_langley(ds, date = f'{which} - {date}')
        return aa

class CalibrationsOverTime(object):
    def __init__(self, path2historic = '/mnt/telg/projects/sp02/calibration/BRWsp02calhist20062018SEDcorr.dat', list0path2modern = [], additional_result_instance = None):
        self.paths = dict(path2historic = pathlib.Path(path2historic),
                          list0path2modern = [pathlib.Path(p) for p in list0path2modern])
        self._historic = None
        self._results = None
        self._modern  = None
        self._additional = None
        self.additional = additional_result_instance
#         self._additional = additional_result_instance
        
    @property
    def additional(self):
        return self._additional

    @additional.setter
    def additional(self,value):
        if not isinstance(self._additional, type(None)):
            self._remove_additional()
        self._additional = value
        self._add_additional()
        
    def _add_additional(self):
        if isinstance(self._additional, type(None)):
            return
        if not isinstance(self._results, type(None)):
            meana, stda = self._ds2resutltdf(self.additional.langley_fitres)
            adt = pd.to_datetime(self.additional.langley_fitres.datetime.mean().values)
            self._results['mean'].loc[adt] = meana
            self._results['std'].loc[adt] = stda
    
    def _remove_additional(self):
        mean = self.results['mean']
        adt = pd.to_datetime(self.additional.langley_fitres.datetime.mean().values)
        mean.drop([adt], inplace = True)
        return 
    
    @property
    def modern(self):
        if isinstance(self._modern, type(None)):
            self._load_modern()
        return self._modern
    
    def _load_modern(self):
        return [] 
    

    @property
    def results(self):
        if isinstance(self._results, type(None)):
            mean = self.historic_cals.copy()
            std = self._historic_std

            self._results = dict(mean = mean,
                                 std = std)
            
            
            for p2m in self.paths['list0path2modern']:
                langley_fitres = xr.open_dataset(p2m)
                meana, stda = self._ds2resutltdf(langley_fitres)
                adt = pd.to_datetime(langley_fitres.datetime.mean().values)
                self._results['mean'].loc[adt] = meana
                self._results['std'].loc[adt] = stda
            
            self._add_additional()
            
            # hist
            # modern
            #additional
        return self._results
            
                      
    @property
    def historic_cals(self):
        if isinstance(self._historic, type(None)):
            self._read_historic()
        return self._historic
    
    def _read_historic(self):
        path2file = self.paths['path2historic']
        df = pd.read_csv(path2file, delim_whitespace=True, skiprows=2,
                    names = ['year', 'doy', 'ch', 'cal'],
                   )

        df.index = pd.to_datetime(df.year, format='%Y') + pd.to_timedelta(df.doy, 'd')

        channels = pd.Series([int(i) for i in '368 1050 610 778 412 500 675 862'.split()])
        channels.index += 1

        df['wl'] = df.ch.apply(lambda x: channels.loc[x])
        df = df.drop(['year', 'doy', 'ch'], axis=1)

        trans_list = []
        for dt in df.index.unique():
            df_tt = df[df.index == dt]
        #     print(df_tt.shape)
            # df_trans.columns = df_trans.loc['wl']
            df_tt = pd.DataFrame([df_tt.cal.values.transpose()], columns=df_tt.wl, index = [dt])
            df_tt.index.name = 'date'
            trans_list.append(df_tt)

        df_v0_ts = pd.concat(trans_list)
        self._historic = df_v0_ts
        self._historic_std = df_v0_ts.copy()
        self._historic_std[:] = np.nan

    def _ds2resutltdf(self,ds):
        out = {}
        dst = np.e**ds.langley_fitres.sel(linreg_param = 'intercept')
        
        cw = ds.channle_wavelengths
        if cw.dims == ('channel',):
            cw = cw.to_pandas()
        elif cw.dims == ('channel','datetime',):
            cw = cw.transpose('channel','datetime').to_pandas().iloc[:,0]
        else:
            raise ValueError(f'should not be possible ....ds.dims = {ds.dims}')

        mean = dst.mean(dim = 'datetime').to_pandas()
    #     print(mean)
    #     print(cw)
        mean.rename(cw, inplace=True)

        std = dst.std(dim = 'datetime').to_pandas()
        std.rename(cw, inplace=True)
        return mean, std
    
    def plot(self, slides = False, serial_no = None):
        out = {}
        df_v0_ts = self.results['mean'].sort_index()
        df_v0_ts_error = self.results['std'].sort_index()
        newcols = df_v0_ts.iloc[-1].dropna().index
        df_v0_ts =df_v0_ts.reindex(newcols, axis=1)
        df_v0_ts_error = df_v0_ts_error.reindex(newcols, axis=1)

        if slides:
            date = pd.to_datetime(datetime.datetime.now()).date()
            date = str(date)
            date = date.replace('-','')
            # path2ppt = f'/mnt/telg/projects/sp02/calibration/mlo_v0_over_time_{date}.ppt'
            path2ppt = f'calibration_{date}.ppt'
            
            if type(slides).__name__ == 'Presentation':
                prs = slides
            else:
                prs = pptx.Presentation()
                if 1:
                    title_slide_layout = prs.slide_layouts[0]
                    slide = prs.slides.add_slide(title_slide_layout)
                    title = slide.shapes.title
                    subtitle = slide.placeholders[1]
        
                    title.text = "v0 over time"
                    now = pd.to_datetime(datetime.datetime.now())
        
                    subtitle.text = f"created on {now.date()} by Hagen"


            ttop = 0
            tleft = 0.1
            lefts = [0.1,5]*2
            tops = [0.2]*2 + [3.5]*2
            width = 4.5
        # average over everything before 18
        ### average over all but last

        oldcal = df_v0_ts.iloc[:-1, :].mean()

        for e,ch in enumerate(df_v0_ts):
            if slides:
                if np.mod(e,4) == 0:
                    blank_slide_layout = prs.slide_layouts[6]
                    slide = prs.slides.add_slide(blank_slide_layout)
                    txBox = slide.shapes.add_textbox(tleft, ttop, width, 0.2)
                    tf = txBox.text_frame
                    tf.text = f'sn #{serial_no}'
            f,a = plt.subplots()
            f.set_size_inches(4.5, 3.5)
        #     df_v0_ts[ch].plot(ax = a)
            a.plot(df_v0_ts[ch].index, df_v0_ts[ch].values,marker = 'o', markersize = 5, ls = '--')

            a.errorbar(df_v0_ts[ch].index, df_v0_ts[ch].values, df_v0_ts_error[ch].values, capsize = 2, ls = '', ecolor = 'red', zorder = 10)
        #     a.errorbar()

            # average over everything before 18
            a.axhline(y = oldcal.loc[ch], color = colors[1])


        #     anhere = pd.to_datetime('20140101')
            anhere = df_v0_ts.index[-1] - pd.to_timedelta(4*365, 'd')
            avg = oldcal.loc[ch]
            last = df_v0_ts[ch].values[-1]
            mid = (avg + last)/2
            diff = (last - avg)/avg *100

            a.annotate('',
                    xy=(anhere, avg), xycoords='data',
                    xytext=(anhere, last), textcoords='data',
                    arrowprops=dict(arrowstyle="<->",
                                    connectionstyle="arc3"),
                    ha = 'center',
                    va = 'center'
                    )
            a.text(anhere, mid, f" {diff:0.1f} %")

            # annotate change from last time
            anhere = df_v0_ts.index[-1] - pd.to_timedelta(1*365, 'd')
            avg = df_v0_ts[ch].values[-2]
            last = df_v0_ts[ch].values[-1]
            mid = (avg + last)/2
            diff = (last - avg)/avg *100
            dist = avg - last

            a.annotate('',
                    xy=(anhere, avg), xycoords='data',
                    xytext=(anhere, last), textcoords='data',
                    arrowprops=dict(arrowstyle="<->",
                                    connectionstyle="arc3"),
                    ha = 'center',
                    va = 'center'
                    )
            a.text(anhere, mid, f" {diff:0.1f} %")



            a.set_title(f'channel: {ch} nm')
            a.grid()
            a.set_ylabel('V0')
            a.set_xlabel('Date of calibration')
            if slides:
                add_fig2slide(f, slide,left = lefts[np.mod(e,4)], top = tops[np.mod(e,4)], width = width)
        #     break
        if slides:
            prs.save(path2ppt)
            out['slides'] = prs
        return out
    


class SP02RawData(object):
    def __init__(self, dataset, site, langley_fit_settings = None):
        self.raw_data = dataset
        self.site = site
        self.langley_fit_settings = langley_fit_settings
        self._sun_position = None
        self._am = None
        self._pm = None
        # self._langleys_am = None
        # self._langleys_pm = None
        # self._langley_fitres_am = None
        # self._langley_fitres_pm = None
        
    @property
    def sun_position(self):
        if isinstance(self._sun_position, type(None)):
            self._sun_position = self.site.get_sun_position(self.raw_data.datetime)
        return self._sun_position
    
    @property
    def am(self):
        if isinstance(self._am, type(None)):
            self._get_langley_from_raw() 
        return self._am
    
    @property
    def pm(self):
        if isinstance(self._pm, type(None)):
            self._get_langley_from_raw() 
        return self._pm
    
        
    
    def _get_langley_from_raw(self):
        raw_df = self.raw_data.raw_data.to_pandas()
        
        # changing to local time
        raw_df_loc = raw_df.copy()
        index_local = raw_df.index + pd.to_timedelta(self.site.time_zone[1], 'h')
        raw_df_loc.index = index_local

        # getting the one day
        start = raw_df_loc.index[0]
        start = pd.to_datetime(f'{start.year}{start.month:02d}{start.day:02d}') + pd.to_timedelta(1,'d')
        end = start + pd.to_timedelta(1, 'd')
        raw_df_loc = raw_df_loc.truncate(start, end)

        # localize and cut day for sunposition
        sunpos = self.sun_position.copy()
        sunpos.index = index_local
        sunpos = sunpos.truncate(start, end)

        # remove the night
        sunpos[sunpos.airmass < 0] = np.nan

        # get the minimum airmass befor I start cutting it out
        noon = sunpos.airmass.idxmin()

        # the data needs to be normalized to the sun_earth_distance
#         self.tp_raw_df_loc = raw_df_loc.copy()
#         self.tp_sunpos = sunpos.copy()
        raw_df_loc = raw_df_loc.multiply(sunpos.sun_earth_distance**2, axis=0)
    
        # langleys are the natural logarith of the voltage over the AMF ... -> log
        # to avoid warnings and strange values do some cleaning before log
        raw_df_loc[raw_df_loc <= 0] = np.nan
#         self.tp_raw_df = raw_df.copy()
        raw_df_loc = np.log(raw_df_loc)    
    
        # keep only what is considered relevant airmasses
        amf_min = 2.2 
        amf_max = 4.7
        sunpos[sunpos.airmass < amf_min] = np.nan
        sunpos[sunpos.airmass > amf_max] = np.nan

        sunpos_am = sunpos.copy()
        sunpos_pm = sunpos.copy()

        sunpos_am[sunpos.index > noon] = np.nan
        sunpos_pm[sunpos.index < noon] = np.nan
        
        langley_am = raw_df_loc.copy()
        langley_pm = raw_df_loc.copy()

        langley_am.index = sunpos_am.airmass
        langley_pm.index = sunpos_pm.airmass

        self._am = Langley(self,langley_am[~langley_am.index.isna()], langley_fit_settings = self.langley_fit_settings)
        self._pm = Langley(self,langley_pm[~langley_pm.index.isna()], langley_fit_settings = self.langley_fit_settings)
        return True

   
    
def raw2langleys(site,
                 start_date = '20200205',
                 end_date = '20200207',
                 path2netcdf = '/mnt/telg/data/baseline/mlo/2020/',
                 pattern = '*.nc',
                 pathout_fld = '.',
                 break_when_error = False,
                 test = False):
    out = {}
    mlo = site
    path2netcdf = pathlib.Path(path2netcdf)
    # path2netcdf.mkdir()

    ### generate folder names and work plan
    runname = 'langleys'
    
    pathout_fld = pathlib.Path(pathout_fld)
    pathout_fld = pathout_fld.resolve()
    pathout_fld.mkdir(exist_ok=True)
    
    pathout_fld_run = pathout_fld.joinpath(runname)
    pathout_fld_run.mkdir(exist_ok=True)
    


    df = pd.DataFrame(list(path2netcdf.glob(pattern)), columns=['path'])
    
    ### start time index
    df.index = df.apply(lambda row: pd.to_datetime(row.path.name.split('.')[4]), axis=1)
    df.sort_index(inplace=True)
    df_files = df
    

    ### return df_files

    df_sel = df.truncate(before = start_date,
                         after = end_date
                        )
    
    ### open first file to get serial number for output names and testing
    ds = xr.open_dataset(df_sel.iloc[0,0])
    serial_no = int(ds.serial_no.values)
    ds.close()
    
    ### generate output files and check if exists
    df_sel['path_out'] = df_sel.apply(lambda x: pathout_fld_run.joinpath(f'sn{serial_no}_{x.name.year:04d}{x.name.month:02d}{x.name.day:02d}_am.nc'), axis = 1)
    df_sel['output_file_exists'] = df_sel.path_out.apply(lambda x: x.is_file())    
    total_no = df_sel.shape[0] - 1 # -1 since I alsways process 2 at a time
    df_sel = df_sel[~df_sel.output_file_exists]
    work_no = df_sel.shape[0] - 1
    

    print(f'processing {work_no} of a total of {total_no} files: ', end = '')
    # return df_sel, serial_no, ds, pathout_fld_run
    
    reslist_am = []
    dstl = []
    # print(f'processing {df_sel.shape[0]} files', end = ' ... ')
    for i in range(len(df_sel)-1):

    #     if i == 2:
    #         break
    #     if i <= 46:
    #         continue
        print(i,end = '.')
        row_duo = df_sel.iloc[i:i+2]
        ds = xr.open_mfdataset(row_duo.path, combine='by_coords')
        out['tp1'] = row_duo.path
        raw = SP02RawData(ds, mlo)
        serial_no_raw = raw.raw_data.serial_no.values[0]
        if serial_no_raw != serial_no:
            raise ValueError('serial no is {serial_no_raw}, but should be {serial_no}')
    #     dstl.append(raw)
    #     resdict_am[row_duo.index[0]] = raw.langley_fitres_am.copy()
        try:
            raw.am
        except:
            print(f'problem with langley retieval at i = {i}')
            if break_when_error:
                break
            else:
                work_no -= 1
                continue


        try:
            raw.am.langley_fitres
        except:
            print(f'problem with fit at i = {i}')
            if break_when_error:
                break
            else:
                work_no -= 1
                continue

        reslist_am.append(dict(datetime = df_sel.index[i],
                               langley_fitres = raw.am.langley_fitres.copy(),
                               langleys = raw.am.langleys.copy(),
                               correlation_prop = raw.am.langley_residual_correlation_prop.copy(),
                               langley_fit_resid = raw.am.langley_fit_residual.copy(),
                              )
                         )
        # datestr = df_sel.index[i].date().__str__().replace('-', '')
        # sn = int(np.unique(raw.raw_data.serial_no.values))
        # pathout_langleyres = pathout_fld_run.joinpath(f'sn{sn}_{datestr}_am.nc')
        
        # raw.am['channle_wavelengths'] = raw.raw_data.channle_wavelengths.to_pandas().iloc[1]
        # return raw
        raw.am.save2netcdf(df_sel.iloc[i].path_out)
        ds.close()
        if test:
            break
    #     break
    print('done')
    
    pathout_file = pathout_fld_run.joinpath(f'sn{serial_no}_all_fitres.nc')
    
    if work_no > 0:
        ### integrate everthing in one Dataset
        fitresarray = np.array([res['langley_fitres'].values for res in reslist_am])
        coord_time = [res['datetime'] for res in reslist_am]
        correlation_props = [res['correlation_prop'] for res in reslist_am]    
        coord_linreg = raw.am.langley_fitres.columns.values
        coord_amf = raw.am.langleys.index
        coord_wavelengths = raw.am.langley_fitres.index.values
        langley_fitres_da = xr.DataArray(fitresarray, coords=[coord_time, coord_wavelengths, coord_linreg], dims=['datetime', 'wavelength', 'linreg_param'])
        
        
        ds = xr.Dataset()
        ds['langley_fitres'] = langley_fitres_da
        # ds['langleys'] = langleys_da
        # ds['langley_fit_residual'] = langleys_resid_da
        
        cmp = pd.DataFrame(correlation_props, index = coord_time)
        cmp.index.name = 'datetime'
        cmp.columns.name = 'matrix_properties'
        ds['correlation_matrix_properties'] = cmp
        
        ds['channle_wavelengths'] = raw.raw_data.channle_wavelengths.to_pandas().iloc[1]
        ds['sp02_serial_no'] = serial_no
    
    
        
        if pathout_file.exists():
            ds_others = xr.open_dataset(pathout_file) 
            ds = xr.concat([ds, ds_others], dim = 'datetime')
            ds = ds.sortby('datetime')
            ds_others.close()
    
        ds.to_netcdf(pathout_file)
    else:
        if pathout_file.exists():
            ds = xr.open_dataset(pathout_file) 
        else:
            ds = None
        raw = None
        
        

    out['ds_results'] = ds
    out['last_raw_instance'] = raw
    return Calibration(out)

def add_fig2slide(f, slide, left = 0, top = 0, width = 5.5, add_text_box = False):
    imstream = io.BytesIO()
    # f.patch.set_alpha(0)
    f.tight_layout()
    f.savefig(imstream, format='png', bbox_inche = 'tight')

#     blank_slide_layout = prs.slide_layouts[6]
#     slide = prs.slides.add_slide(blank_slide_layout)

    left =  pptx.util.Inches(left)
    top = pptx.util.Inches(top)
    width = pptx.util.Inches(width)
    pic = slide.shapes.add_picture(imstream, left, top, width = width)


#         slide = prs.slides.add_slide(blank_slide_layout)
    left = width # pptx.util.Inches(0.5)
#         width = pptx
    top = top + pptx.util.Inches(0.3)
    height = pptx.util.Inches(3)
    width = pptx.util.Inches(4)
    if add_text_box:
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = '...'
    return

def deprecated_plot_langley_results(result_ds, intensity4ts = 'resid_curvature', add2pptx = False,test = False, tollerance = [2.5, 2.5]):
    """
    slope 	intercept 	slope_stderr 	intercept_stderr 	resid_curvature

    Parameters
    ----------
    result_ds : TYPE
        DESCRIPTION.
    intensity4ts : TYPE, optional
        DESCRIPTION. The default is ''.
    add2pptx : TYPE, optional
        DESCRIPTION. The default is False.
    test : TYPE, optional
        DESCRIPTION. The default is False.
    tollerance : TYPE, optional
        DESCRIPTION. The default is [2.5, 2.5].

    Returns
    -------
    out : TYPE
        DESCRIPTION.

    """
    ds = result_ds
    out = {}
    if add2pptx:
        date = pd.to_datetime(datetime.datetime.now()).date()
        date = str(date)
        date = date.replace('-','')
        path2ppt = f'/mnt/telg/projects/sp02/calibration/langley_calibration_summary_{date}.ppt'
        path2ppt
        prs = pptx.Presentation()
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = "SP02 revival"
        now = pd.to_datetime(datetime.datetime.now())

        subtitle.text = f"created on {now.date()} by Hagen"



    plt.rcParams['hatch.linewidth'] = 4
    for wl in ds.wavelength.values:
        f,a = plt.subplots()
        res = ds.langley_fitres.sel(wavelength = wl).to_pandas()

    #     x,y,z = res.slope.copy(), res.stderr.copy(), res.intercept.copy()
        rest = res.copy()
        if 0:
            c = rest.slope.median()
            d = abs(c) * 0.5
            xlim =((c - d),(c + d))
            rest.slope[rest.slope > xlim[1]] = np.nan
            rest.slope[rest.slope < xlim[0]] = np.nan
        else:
            xlim = (None, None)
            
        
        if 1:
            bottom = res.intercept_stderr.min()
            top = tollerance[1] * bottom
            dlim = (top - bottom) * 0.1
            ylim = (bottom - dlim , top + dlim)
    #         print(ylim)
        #     ylim =(0,(c + d))
            rest.intercept_stderr[rest.intercept_stderr > top] = np.nan
            rest.intercept_stderr[rest.intercept_stderr < bottom] = np.nan
        else:
            ylim = (None, None)

        restdna = rest.dropna()

        noall = res.shape[0]
        inside = restdna.shape[0]
    #     inside = np.logical_and(~x.isna(), ~y.isna()).sum()
        outside = noall - inside
        txt = f'{outside} of {noall} data\npoints not within\nplotting limits'
    #     print(txt)
        a.text(0.95,0.95, txt,ha = 'right', va = 'top', transform=a.transAxes,)
    #     mi = 0
    #     ma = 10
    #     y[y > ma] = np.nan
    #     y[y < mi] = np.nan

        if 0:
            pc = plt.hexbin(res.slope, res.stderr, gridsize=100)
            cb = f.colorbar(pc)
        #     pc.set_linewidth(0.2)
            pc.set_edgecolor([0,0,0,0])
            cm = plt.cm.plasma_r
            pc.set_cmap(cm)
            cm.set_under('w')
            pc.set_clim(0.01)
            cb.set_label('# of datapoints on bin')

        elif 0:
            g, = plt.plot(x, y)
            g.set_linestyle('')
            g.set_marker('.')
            g.set_markersize(1)
            a.set_xlim(xlim)
            a.set_ylim(ylim)
        else:
            cm = plt.cm.plasma_r
            pc = plt.scatter(restdna.slope,restdna.intercept_stderr, c = restdna.intercept, s = 40, cmap = cm, 
    #                     edgecolors='black',
                       )
            cb = f.colorbar(pc)
            cb.set_label('Langley fit - intercept')
            zm = restdna.intercept.median()
            zstd  = restdna.intercept.std()
            cb.ax.axhline(zm, color = 'black')
            hs = cb.ax.axhspan(zm-zstd, zm+zstd)
            hs.set_facecolor([0,0,0,0])
            hs.set_edgecolor([0,0,0,0.5])
    # #         hs.set_edgecolor([1,1,1,1])
            hatch = hs.set_hatch('//')
    #         hs.set_linewidth(10)
            a.set_xlim(xlim)
            a.set_ylim(ylim)
    #         print(ylim)
        a.set_title(f'channle wavelength: {wl}nm')
        a.set_xlabel('Langley fit - slope (mV)')
        a.set_ylabel('Langley fit - std of residual')

    # langley over time
        f_ivt,a = plt.subplots()
        out['tp_restdna'] = restdna
        # restdna.intercept.plot(ax = a)
        pc = a.scatter(restdna.index, restdna.intercept, c = restdna[intensity4ts], cmap=plt.cm.gnuplot)
        f_ivt.colorbar(pc)
        # g = a.get_lines()[0]
        # g.set_linestyle('')
        # g.set_marker('o')
        a.set_xlabel('')
        a.set_ylabel('Langley fit - intercept')
        a.set_title('Langley fit intercept as a function of time')

        if add2pptx:

            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            add_fig2slide(f, slide)
            add_fig2slide(f_ivt, slide, top = 7.5/2)
        if test:
            break
    if add2pptx:
        prs.save(path2ppt)
    return out

def deprecated_get_best_results(ds, top = 10):
    def sortbyparam(ds, param, top):
        df = (ds.langley_fitres.sel(linreg_param = f'{param}_stderr') / ds.langley_fitres.sel(linreg_param = param)).mean(dim = 'wavelength').to_pandas()
        df = df[df>0]
        df.sort_values(inplace=True)
        out = {}
        df = df.iloc[:top]
        out = ds.sel(datetime =  df.index.values)
        return out
    out = {}
    slopestdrmin = sortbyparam(ds,'slope', top)
    interceptstdrmin = sortbyparam(ds,'intercept',top)
    curvature_date = ds.langley_fitres.sel(linreg_param = 'resid_curvature').mean(dim = 'wavelength').to_pandas().abs().sort_values().index[:top].values
    curvature = ds.sel(datetime = curvature_date)
    df_cmp = ds.correlation_matrix_properties.to_pandas()
    determimin_dates = df_cmp.sort_values('determinant', ascending=False).index[:top].values
    determimin = ds.sel(datetime = determimin_dates)
#     uberlap = ds.datetime.values
#     for other in [interceptstdrmin, 
#                  curvature, 
#     #              determimin,
#                  ]:
#         uberlap = np.intersect1d(uberlap, other)
#     out['crossection'] = uberlap
    
    out['slope_stderr'] = slopestdrmin
    out['intercept_stderr'] = interceptstdrmin
    out['curvature'] = curvature
    out['corr_matrix_determinant_max'] = determimin
    return out


def future_add_fig2slide(f, slide, left = 0, top = 0, width = 5.5):
    imstream = io.BytesIO()
    f.patch.set_alpha(0)
    f.tight_layout()
    f.savefig(imstream, format='png', bbox_inche = 'tight')

#     blank_slide_layout = prs.slide_layouts[6]
#     slide = prs.slides.add_slide(blank_slide_layout)

    left =  pptx.util.Inches(left)
    top = pptx.util.Inches(top)
    width = pptx.util.Inches(width)
    pic = slide.shapes.add_picture(imstream, left, top, width = width)


#         slide = prs.slides.add_slide(blank_slide_layout)
#     left = width # pptx.util.Inches(0.5)
# #         width = pptx
#     top = top + pptx.util.Inches(0.3)
#     height = pptx.util.Inches(3)
#     width = pptx.util.Inches(4)
#     txBox = slide.shapes.add_textbox(left, top, width, height)
#     tf = txBox.text_frame
#     tf.text = '...'
    return